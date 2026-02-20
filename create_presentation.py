#!/usr/bin/env python3
"""Create a 3-slide Resilience consulting presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
DARK_NAVY = RGBColor(0x0B, 0x1D, 0x3A)
ACCENT_BLUE = RGBColor(0x1B, 0x5E, 0x9E)
LIGHT_BLUE = RGBColor(0x3A, 0x8F, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MID_GRAY = RGBColor(0x6B, 0x7B, 0x8D)
HIGHLIGHT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
HIGHLIGHT_GREEN = RGBColor(0x00, 0x8A, 0x5E)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)


def add_bg(slide, color):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a filled rectangle shape as a background element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from pptx.oxml.ns import qn
        solidFill = shape.fill._fill
        srgbClr = solidFill.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        alphaElem = srgbClr.makeelement(qn("a:alpha"), {"val": str(alpha)})
        srgbClr.append(alphaElem)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.15):
    """Add a text box with a single run of text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=DARK_TEXT, bullet_color=ACCENT_BLUE, font_name="Calibri",
                    spacing=6, bold_prefix=True):
    """Add a bulleted text list. Items can be 'Bold Part | rest' or plain text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.space_before = Pt(2)
        p.level = 0

        # Bullet character
        bullet_run = p.add_run()
        bullet_run.text = "\u25cf  "
        bullet_run.font.size = Pt(font_size - 2)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.name = font_name

        if bold_prefix and "|" in item:
            bold_part, rest = item.split("|", 1)
            run_bold = p.add_run()
            run_bold.text = bold_part
            run_bold.font.size = Pt(font_size)
            run_bold.font.bold = True
            run_bold.font.color.rgb = color
            run_bold.font.name = font_name

            run_rest = p.add_run()
            run_rest.text = rest
            run_rest.font.size = Pt(font_size)
            run_rest.font.bold = False
            run_rest.font.color.rgb = color
            run_rest.font.name = font_name
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.bold = False
            run.font.color.rgb = color
            run.font.name = font_name

    return txBox


def add_card(slide, left, top, width, height, title, body_items,
             accent_color=ACCENT_BLUE, font_size=11):
    """Add a card-style box with title bar and bullet content."""
    # Card background
    card_bg = add_shape_bg(slide, left, top, width, height, WHITE)
    # Accent top bar
    add_shape_bg(slide, left, top, width, Inches(0.06), accent_color)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4),
                 Inches(0.35), title, font_size=13, bold=True, color=accent_color)
    # Body bullets
    add_bullet_list(slide, left + Inches(0.15), top + Inches(0.48),
                    width - Inches(0.3), height - Inches(0.6),
                    body_items, font_size=font_size, color=DARK_TEXT,
                    bullet_color=accent_color, spacing=4, bold_prefix=True)


def add_stat_box(slide, left, top, width, height, number, label,
                 num_color=ACCENT_BLUE, label_color=MID_GRAY):
    """Add a statistic highlight box."""
    box = add_shape_bg(slide, left, top, width, height, WHITE)
    add_text_box(slide, left, top + Inches(0.15), width, Inches(0.5),
                 number, font_size=26, bold=True, color=num_color,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.65), width - Inches(0.2),
                 Inches(0.6), label, font_size=10, bold=False, color=label_color,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — THE RESILIENCE MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide1, LIGHT_GRAY)

# Left accent panel
add_shape_bg(slide1, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT_BLUE)

# Header bar
add_shape_bg(slide1, Inches(0.12), Inches(0), Inches(13.213), Inches(1.1), DARK_NAVY)
add_text_box(slide1, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
             "THE RESILIENCE MARKET OPPORTUNITY",
             font_size=28, bold=True, color=WHITE, font_name="Calibri")
add_text_box(slide1, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35),
             "Critical Infrastructure | Business Continuity | Threat Monitoring | Asset Protection",
             font_size=13, bold=False, color=LIGHT_BLUE, font_name="Calibri")

# ── Row of stat boxes ──
stat_y = Inches(1.35)
stats = [
    ("~$150B", "Critical Infrastructure\nProtection Market (2025)"),
    ("$2\u2013$3B", "Business Continuity\nManagement Market (2025)"),
    ("8\u201317%", "BCM Market CAGR\n(2025\u20132033)"),
    ("5.8%", "Services Segment\nCAGR (fastest growing)"),
    ("$1.06T", "Global Consulting\nMarket (2025)"),
]
stat_width = Inches(2.2)
stat_gap = Inches(0.27)
start_x = Inches(0.6)
for i, (num, label) in enumerate(stats):
    x = start_x + i * (stat_width + stat_gap)
    add_stat_box(slide1, x, stat_y, stat_width, Inches(1.2), num, label)

# ── Left column: Market drivers ──
col1_x = Inches(0.6)
col1_y = Inches(2.85)
add_text_box(slide1, col1_x, col1_y, Inches(6), Inches(0.35),
             "KEY MARKET DRIVERS", font_size=15, bold=True, color=DARK_NAVY)

drivers = [
    "Regulatory Surge:|  EU DORA (live Jan 2025), NIS2 (deadline Oct 2026), CER (entity ID by Jul 2026) creating multi-framework compliance demand across all sectors",
    "Escalating Threat Landscape:|  State-backed threats (esp. PRC), ransomware breakout times <60 min, identity-based attacks surpassing malware as primary intrusion vector",
    "AI Dual-Edge:|  AI accelerating both attack sophistication and defensive resilience capabilities; AI supply chain attacks emerging as a new attack surface",
    "Supply Chain Fragility:|  High-profile incidents (JLR, M&S) exposing third-party resilience gaps; regulators mandating vendor risk management integration",
    "Digital Transformation:|  Multi-cloud resilience, zero-trust infrastructure, and real-time risk monitoring now table stakes for enterprise operations",
]
add_bullet_list(slide1, col1_x, col1_y + Inches(0.4), Inches(6.2), Inches(3.5),
                drivers, font_size=12, color=DARK_TEXT, bullet_color=HIGHLIGHT_ORANGE,
                spacing=6)

# ── Right column: Target Sectors ──
col2_x = Inches(7.2)
add_text_box(slide1, col2_x, col1_y, Inches(5.5), Inches(0.35),
             "HIGH-DEMAND SECTORS", font_size=15, bold=True, color=DARK_NAVY)

sectors = [
    "Financial Services:|  DORA compliance, operational resilience frameworks, third-party risk management, ICT incident reporting",
    "Energy & Utilities:|  28.9% of CIP market; grid resilience, OT security, green transition risk management",
    "Healthcare:|  Digital health resilience, medical device security, ransomware recovery, regulatory compliance",
    "Government & Defence:|  Federal continuity programs, CISA coordination, critical infrastructure protection planning",
    "Transport & Supply Chain:|  Growing at 4.7% CAGR to 2031; logistics resilience, port/aviation security",
]
add_bullet_list(slide1, col2_x, col1_y + Inches(0.4), Inches(5.8), Inches(3.5),
                sectors, font_size=12, color=DARK_TEXT, bullet_color=HIGHLIGHT_GREEN,
                spacing=6)

# Bottom bar with verdict
add_shape_bg(slide1, Inches(0.12), Inches(6.8), Inches(13.213), Inches(0.7), DARK_NAVY)
add_text_box(slide1, Inches(0.6), Inches(6.88), Inches(12.5), Inches(0.5),
             "VERDICT:  Strong and growing market. Regulatory deadlines in 2025\u20132026 are creating immediate, sustained demand for resilience consulting services across industries.",
             font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — RECENT TRENDS (Dec 2025 \u2013 Feb 2026)
# ═══════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, LIGHT_GRAY)
add_shape_bg(slide2, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT_BLUE)

# Header
add_shape_bg(slide2, Inches(0.12), Inches(0), Inches(13.213), Inches(1.1), DARK_NAVY)
add_text_box(slide2, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
             "RESILIENCE TRENDS  |  Dec 2025 \u2013 Feb 2026",
             font_size=28, bold=True, color=WHITE)
add_text_box(slide2, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35),
             "What\u2019s happening in hiring, regulation, thought leadership, and technology",
             font_size=13, bold=False, color=LIGHT_BLUE)

# ── Card Row 1 ──
card_w = Inches(3.95)
card_h = Inches(2.85)
row1_y = Inches(1.35)
gap = Inches(0.3)

add_card(slide2, Inches(0.6), row1_y, card_w, card_h,
         "REGULATORY MOMENTUM",
         [
             "DORA:|  Fully in force since Jan 2025; Register of Information submitted Apr 2025; critical ICT vendor assessments in progress",
             "NIS2:|  14 of 27 EU states transposed; compliance deadline Oct 2026; Germany, Portugal, Austria recently adopted",
             "CER:|  Member states must identify critical entities by Jul 2026; covers 11 key sectors",
             "UK:|  New Cyber Security & Resilience Bill aligning with EU standards",
         ], accent_color=ACCENT_BLUE, font_size=11)

add_card(slide2, Inches(0.6) + card_w + gap, row1_y, card_w, card_h,
         "THREAT LANDSCAPE",
         [
             "State-backed threats:|  PRC pre-positioning across US critical infrastructure (CISA priority)",
             "Ransomware:|  Breakout times under 60 minutes; credential abuse now surpasses malware as primary vector",
             "AI-powered attacks:|  Threat actors using LLMs for phishing, recon, malware dev; AI supply chain attacks emerging",
             "Incidents:|  JLR and M&S disruptions underscored supply chain fragility; food/ag sector targeted",
         ], accent_color=RGBColor(0xC0, 0x39, 0x2B), font_size=11)

add_card(slide2, Inches(0.6) + 2 * (card_w + gap), row1_y, card_w, card_h,
         "HIRING & WORKFORCE",
         [
             "39%| of organisations hiring more digital trust roles in 2026 vs 2025",
             "Nearly 50%| expect difficulty finding qualified candidates in audit, risk, and cybersecurity",
             "Skills > headcount:|  Organisations prioritising skill-based hiring and upskilling over expanding teams",
             "Key gaps:|  Data security, emerging tech risk, AI governance, OT security specialists",
         ], accent_color=HIGHLIGHT_GREEN, font_size=11)

# ── Card Row 2 ──
row2_y = Inches(4.5)
card_h2 = Inches(2.3)

add_card(slide2, Inches(0.6), row2_y, card_w, card_h2,
         "TECHNOLOGY SHIFTS",
         [
             "AI-powered resilience:|  Real-time risk monitoring, predictive infrastructure planning, automated incident response",
             "Zero Trust:|  Embedded directly into infrastructure management; no longer a separate security layer",
             "Multi-cloud resilience:|  Enterprises re-architecting for multi-region failover and uptime SLAs",
         ], accent_color=RGBColor(0x7D, 0x3C, 0x98), font_size=11)

add_card(slide2, Inches(0.6) + card_w + gap, row2_y, card_w, card_h2,
         "THOUGHT LEADERSHIP & WHITEPAPERS",
         [
             "ISACA:|  \"Resilience and Security in Critical Sectors: Navigating NIS2 & DORA\" (2025)",
             "CISA:|  Infrastructure Resilience Planning Framework updated Jan 2025; 2025 Year in Review (Feb 2026)",
             "Food & Ag-ISAC:|  \"Farm-to-Table Ransomware Realities \u2013 2025 Threat Trends & 2026 Outlook\"",
         ], accent_color=HIGHLIGHT_ORANGE, font_size=11)

add_card(slide2, Inches(0.6) + 2 * (card_w + gap), row2_y, card_w, card_h2,
         "MARKET SIGNALS",
         [
             "CISA:|  Stopped 2.62B malicious connections on federal networks in 2025; 148 resilience exercises with 10K+ participants",
             "CIPRNA Expo 2026:|  Dedicated critical infrastructure resilience conference (Mar 2026, Baton Rouge)",
             "Consulting demand:|  Resilience advisory dominates alongside AI adoption as top growth area for 2025\u20132026",
         ], accent_color=ACCENT_BLUE, font_size=11)

# Bottom bar
add_shape_bg(slide2, Inches(0.12), Inches(7.05), Inches(13.213), Inches(0.45), DARK_NAVY)
add_text_box(slide2, Inches(0.6), Inches(7.1), Inches(12.5), Inches(0.35),
             "Sources: CISA, ISACA, Security Magazine, ISC2, Eye Security, Nextgov/FCW, ClearRisk, Hitachi Cyber, Industrial Cyber (Dec 2025 \u2013 Feb 2026)",
             font_size=10, bold=False, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — COMPETITIVE LANDSCAPE: WHAT THE BIG FIRMS ARE DOING
# ═══════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, LIGHT_GRAY)
add_shape_bg(slide3, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT_BLUE)

# Header
add_shape_bg(slide3, Inches(0.12), Inches(0), Inches(13.213), Inches(1.1), DARK_NAVY)
add_text_box(slide3, Inches(0.6), Inches(0.15), Inches(10), Inches(0.55),
             "COMPETITIVE LANDSCAPE: MAJOR CONSULTING FIRMS & RESILIENCE",
             font_size=26, bold=True, color=WHITE)
add_text_box(slide3, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35),
             "How EY, Deloitte, PwC, McKinsey, BCG, and Kearney are positioning in the resilience market",
             font_size=13, bold=False, color=LIGHT_BLUE)

# ── Firm cards ── 3 columns x 2 rows
firm_w = Inches(3.95)
firm_h = Inches(2.55)
firm_gap = Inches(0.3)
firm_row1_y = Inches(1.35)
firm_row2_y = Inches(4.15)

# Deloitte
add_card(slide3, Inches(0.6), firm_row1_y, firm_w, firm_h,
         "DELOITTE",
         [
             "Most extensive| public-facing resilience practice among major firms",
             "Primary consulting provider| for FEMA\u2019s National Continuity Program; shaped Federal Continuity Directives",
             "Dedicated offerings:| OT Resilience, Cyber Resilience, End-to-End Business Resilience, Third-Party Resilience, Crisis & Resilience",
             "Industry focus:| Financial Services, Government/Defence, Healthcare, Energy",
             "2026:| Presenting at CIPRNA Expo on AI agents and predictive infrastructure planning",
         ], accent_color=RGBColor(0x86, 0xBC, 0x25), font_size=10.5)

# EY
add_card(slide3, Inches(0.6) + firm_w + firm_gap, firm_row1_y, firm_w, firm_h,
         "EY",
         [
             "25,000 risk professionals| globally across integrated risk, cyber, forensics, compliance",
             "Operational Resilience platform| (powered by ServiceNow): real-time early warning, automated prevention/recovery",
             "Enterprise Resilience Solution:| holistic approach to reducing risk and driving long-term value through disruption",
             "Industry focus:| Financial Services (regulatory-driven), India risk landscape (FICCI-EY 2026)",
             "Emphasis on:| Regulatory compliance, financial crime, operational resilience, AI-related risks",
         ], accent_color=RGBColor(0xFF, 0xE6, 0x00), font_size=10.5)

# PwC
add_card(slide3, Inches(0.6) + 2 * (firm_w + firm_gap), firm_row1_y, firm_w, firm_h,
         "PwC",
         [
             "\"Resilience by Design\"| philosophy \u2014 shifting from traditional DR/BC to embedded resilience",
             "Cross-industry view:| tailored perspectives for FS, Retail, Healthcare, Tech/Media/Telecom, Manufacturing",
             "Global Centre for Crisis & Resilience| (GCCR) applying Enterprise Resilience Framework",
             "Restructured into 8 advisory platforms| (Jul 2025) inc. Cyber Data & Tech Risk; Risk & Regulatory",
             "Technology & Operational Resilience| practice covers cyber, third-party risk, regulatory compliance",
         ], accent_color=RGBColor(0xD9, 0x3F, 0x0B), font_size=10.5)

# McKinsey
add_card(slide3, Inches(0.6), firm_row2_y, firm_w, firm_h,
         "McKINSEY",
         [
             "Global Infrastructure Initiative (GII):| senior leader community focused on smart, resilient infrastructure",
             "Strategy-level focus:| infrastructure investment resilience, energy transition, grid hardening",
             "Research insight:| nimble networks of smaller units are more resilient than legacy infrastructure",
             "Industry focus:| Energy/Utilities, Transport, Public Sector",
             "15,000+ digital engagements;| 40% YoY growth in sustainability projects",
         ], accent_color=RGBColor(0x00, 0x51, 0x9E), font_size=10.5)

# BCG
add_card(slide3, Inches(0.6) + firm_w + firm_gap, firm_row2_y, firm_w, firm_h,
         "BCG",
         [
             "Operational Risk & Resilience| dedicated practice under Risk Management & Compliance",
             "Key research:| resilient companies achieve 3\u20135% higher annual revenue growth vs peers",
             "Focus areas:| digital transformation, cybersecurity, large-scale transformation, data analytics",
             "Industry coverage:| Healthcare, Industrial Goods, Consumer Products, Energy",
             "Strategic positioning:| resilience as a competitive advantage, not just risk mitigation",
         ], accent_color=RGBColor(0x2B, 0x8C, 0x2B), font_size=10.5)

# Kearney + White Space
add_card(slide3, Inches(0.6) + 2 * (firm_w + firm_gap), firm_row2_y, firm_w, firm_h,
         "KEARNEY + YOUR OPPORTUNITY",
         [
             "Kearney:| global firm with European roots; automotive, healthcare, energy, public sector, transportation",
             "Kearney gap:| no prominently branded standalone resilience practice or published resilience framework",
             "Market white space:| industry-specific resilience services remain fragmented across firms",
             "Opportunity:| niche/boutique firms can compete on speed, specialisation, and regulatory depth (esp. DORA/NIS2/CER)",
             "Differentiator:| end-to-end resilience across cyber + physical + operational domains for mid-market clients",
         ], accent_color=HIGHLIGHT_ORANGE, font_size=10.5)

# Bottom verdict bar
add_shape_bg(slide3, Inches(0.12), Inches(6.95), Inches(13.213), Inches(0.55), DARK_NAVY)
add_text_box(slide3, Inches(0.6), Inches(6.98), Inches(12.5), Inches(0.5),
             "OPPORTUNITY:  Big firms dominate large enterprise. Mid-market and industry-specific resilience consulting (esp. EU regulatory compliance) remains underserved \u2014 a clear opening for specialist consultancies.",
             font_size=13, bold=True, color=WHITE)

# ── Save ──
output_path = "/home/user/Test1/Resilience_Market_Briefing_Feb2026.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
